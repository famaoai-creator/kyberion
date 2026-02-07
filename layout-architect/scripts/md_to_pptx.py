import sys
from pptx import Presentation
from pptx.util import Inches, Pt

def create_pptx(md_path, output_path):
    prs = Presentation()
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Simple splitting by Marp horizontal rule '---'
    slides_content = content.split('---')
    
    for slide_text in slides_content:
        # Create a blank slide with title and content layout
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        lines = [l.strip() for l in slide_text.strip().split('\n') if l.strip()]
        if not lines: continue
        
        # Identify title (lines starting with #)
        title_lines = [l for l in lines if l.startswith('#')]
        body_lines = [l for l in lines if not l.startswith('#') and not l.startswith('<!--')]
        
        if title_lines:
            slide.shapes.title.text = title_lines[0].replace('#', '').strip()
        
        if body_lines:
            tf = slide.placeholders[1].text_frame
            tf.text = body_lines[0].replace('*', '').replace('-', '').strip()
            for line in body_lines[1:]:
                p = tf.add_paragraph()
                p.text = line.replace('*', '').replace('-', '').strip()
                p.level = 0 if not line.startswith((' ', '\t')) else 1

    prs.save(output_path)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 md_to_pptx.py <input.md> <output.pptx>")
        sys.exit(1)
    create_pptx(sys.argv[1], sys.argv[2])
